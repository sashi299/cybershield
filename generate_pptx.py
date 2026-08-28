import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PPTX_PATH = r"C:\Users\hp\Documents\cybershield\CyberShield_Presentation.pptx"

# Color Palette (Dark Cyber Theme)
COLOR_BG = RGBColor(11, 17, 32)         # #0B1120 Dark Navy
COLOR_CARD_BG = RGBColor(30, 41, 59)    # #1E293B Slate Blue
COLOR_CARD_BORDER = RGBColor(51, 65, 85)# #334155
COLOR_PRIMARY = RGBColor(6, 182, 212)   # #06B6D4 Cyan
COLOR_SECONDARY = RGBColor(56, 189, 248)# #38BDF8 Sky Blue
COLOR_SUCCESS = RGBColor(16, 185, 129)  # #10B981 Emerald Green
COLOR_DANGER = RGBColor(239, 68, 68)    # #EF4444 Coral Red
COLOR_WARNING = RGBColor(245, 158, 11)  # #F59E0B Amber
COLOR_TEXT_MAIN = RGBColor(241, 245, 249)# #F1F5F9 White
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# #94A3B8 Gray

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, badge_text, title_text, subtitle_text, slide_num):
    # Top Cyan Accent Line
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.color.rgb = COLOR_PRIMARY

    # Badge Pill
    badge_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.2), Inches(0.35))
    badge_box.fill.solid()
    badge_box.fill.fore_color.rgb = RGBColor(15, 34, 53)
    badge_box.line.color.rgb = COLOR_PRIMARY
    badge_box.line.width = Pt(1)
    tf_b = badge_box.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = badge_text.upper()
    p_b.font.size = Pt(9)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_PRIMARY
    p_b.font.name = "Arial"
    p_b.alignment = PP_ALIGN.CENTER

    # Slide Number
    num_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.4), Inches(2.0), Inches(0.35))
    p_n = num_box.text_frame.paragraphs[0]
    p_n.text = f"SLIDE {slide_num:02d} / 12"
    p_n.font.size = Pt(10)
    p_n.font.bold = True
    p_n.font.color.rgb = COLOR_TEXT_MUTED
    p_n.font.name = "Consolas"
    p_n.alignment = PP_ALIGN.RIGHT

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.6))
    p_t = title_box.text_frame.paragraphs[0]
    p_t.text = title_text
    p_t.font.size = Pt(22)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_TEXT_MAIN
    p_t.font.name = "Arial"

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.4))
    p_s = sub_box.text_frame.paragraphs[0]
    p_s.text = subtitle_text
    p_s.font.size = Pt(11)
    p_s.font.color.rgb = COLOR_TEXT_MUTED
    p_s.font.name = "Arial"

def add_card(slide, left, top, width, height, title, bullets, mode="default"):
    # Background Box
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    
    title_color = COLOR_SECONDARY
    if mode == "danger":
        card.fill.fore_color.rgb = RGBColor(38, 20, 28)
        card.line.color.rgb = COLOR_DANGER
        title_color = COLOR_DANGER
    elif mode == "success":
        card.fill.fore_color.rgb = RGBColor(16, 37, 34)
        card.line.color.rgb = COLOR_SUCCESS
        title_color = COLOR_SUCCESS
    elif mode == "highlight":
        card.fill.fore_color.rgb = RGBColor(18, 38, 58)
        card.line.color.rgb = COLOR_PRIMARY
        title_color = COLOR_PRIMARY
    else:
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

    card.line.width = Pt(1.2)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.18)

    # Title
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color
    p_title.font.name = "Arial"
    p_title.space_after = Pt(8)

    # Bullets
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = f"•  {b}"
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = COLOR_TEXT_MAIN
        p_b.font.name = "Calibri"
        p_b.space_after = Pt(4)

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = f"SPEAKER SCRIPT:\n{notes_text}"

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # completely blank layout

    # =========================================================================
    # SLIDE 1: Title Hero
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)
    
    # Top Accent
    top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.color.rgb = COLOR_PRIMARY

    # Hero Shield Icon
    icon_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.333), Inches(1.0))
    p_icon = icon_box.text_frame.paragraphs[0]
    p_icon.text = "🛡️"
    p_icon.font.size = Pt(54)
    p_icon.alignment = PP_ALIGN.CENTER

    # Main Title
    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.333), Inches(1.1))
    p_t = t_box.text_frame.paragraphs[0]
    p_t.text = "CYBERSHIELD"
    p_t.font.size = Pt(44)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_TEXT_MAIN
    p_t.font.name = "Arial"
    p_t.alignment = PP_ALIGN.CENTER

    # Tagline
    tag_box = s1.shapes.add_textbox(Inches(1.0), Inches(3.4), Inches(11.333), Inches(0.8))
    p_tag = tag_box.text_frame.paragraphs[0]
    p_tag.text = "Real-Time AI-Powered Phishing, Smishing & Fraud Protection Engine"
    p_tag.font.size = Pt(16)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_SECONDARY
    p_tag.font.name = "Arial"
    p_tag.alignment = PP_ALIGN.CENTER

    p_subtag = tag_box.text_frame.add_paragraph()
    p_subtag.text = "Autonomous Lockscreen & In-App Security for Messaging Ecosystems"
    p_subtag.font.size = Pt(13)
    p_subtag.font.color.rgb = COLOR_TEXT_MUTED
    p_subtag.font.name = "Arial"
    p_subtag.alignment = PP_ALIGN.CENTER

    # Meta Chips (3 Columns)
    chip_w = Inches(3.6)
    chip_h = Inches(1.0)
    chip_y = Inches(4.8)

    c1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), chip_y, chip_w, chip_h)
    c1.fill.solid()
    c1.fill.fore_color.rgb = COLOR_CARD_BG
    c1.line.color.rgb = COLOR_PRIMARY
    c1.text_frame.margin_top = Inches(0.15)
    p1 = c1.text_frame.paragraphs[0]
    p1.text = "🚀 TRACK & DOMAIN"
    p1.font.size = Pt(9)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY
    p1.alignment = PP_ALIGN.CENTER
    p1_sub = c1.text_frame.add_paragraph()
    p1_sub.text = "Cybersecurity & AI / ML"
    p1_sub.font.size = Pt(12)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = COLOR_TEXT_MAIN
    p1_sub.alignment = PP_ALIGN.CENTER

    c2 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.86), chip_y, chip_w, chip_h)
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLOR_CARD_BG
    c2.line.color.rgb = COLOR_SUCCESS
    c2.text_frame.margin_top = Inches(0.15)
    p2 = c2.text_frame.paragraphs[0]
    p2.text = "📱 PLATFORM STACK"
    p2.font.size = Pt(9)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_SUCCESS
    p2.alignment = PP_ALIGN.CENTER
    p2_sub = c2.text_frame.add_paragraph()
    p2_sub.text = "Flutter + Native Android + React 19"
    p2_sub.font.size = Pt(12)
    p2_sub.font.bold = True
    p2_sub.font.color.rgb = COLOR_TEXT_MAIN
    p2_sub.alignment = PP_ALIGN.CENTER

    c3 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.73), chip_y, chip_w, chip_h)
    c3.fill.solid()
    c3.fill.fore_color.rgb = COLOR_CARD_BG
    c3.line.color.rgb = COLOR_SECONDARY
    c3.text_frame.margin_top = Inches(0.15)
    p3 = c3.text_frame.paragraphs[0]
    p3.text = "⚡ SPEED & LATENCY"
    p3.font.size = Pt(9)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_SECONDARY
    p3.alignment = PP_ALIGN.CENTER
    p3_sub = c3.text_frame.add_paragraph()
    p3_sub.text = "<10ms Hybrid Core Inference"
    p3_sub.font.size = Pt(12)
    p3_sub.font.bold = True
    p3_sub.font.color.rgb = COLOR_TEXT_MAIN
    p3_sub.alignment = PP_ALIGN.CENTER

    add_speaker_notes(s1, "Good morning judges and mentors. Today, we are proud to present CyberShield — an intelligent, real-time protection system designed to safeguard users from the rapidly growing menace of phishing, WhatsApp smishing, and fraudulent QR codes.")

    # =========================================================================
    # SLIDE 2: Problem Statement
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "Challenge & Market Need", "🚨 The Epidemic of Social Engineering & Smishing", "Why existing cybersecurity tools fail everyday mobile users", 2)

    col_w = Inches(5.66)
    row_h = Inches(2.25)
    add_card(s2, Inches(0.8), Inches(2.0), col_w, row_h, "📱 Smishing & WhatsApp Fraud Surge", [
        "Over 3.4 billion phishing messages are sent globally every day.",
        "Attackers impersonate trusted banks, courier deliveries, electricity bills, and lottery schemes.",
        "Manufactured urgency forces victims into hasty actions ('Account suspended in 2 hours')."
    ], "danger")

    add_card(s2, Inches(6.86), Inches(2.0), col_w, row_h, "📷 The Rise of Quishing (QR Phishing)", [
        "Fraudulent QR codes placed over legitimate merchant payment stands and public posters.",
        "Bypasses standard email/browser gateways since QR codes conceal the true underlying URL from the human eye."
    ], "danger")

    add_card(s2, Inches(0.8), Inches(4.5), col_w, row_h, "⏳ The Failure of Manual Security Checkers", [
        "Existing scanners require users to manually copy-paste links into a browser website.",
        "Victims click first due to urgency or trust — by then, credentials or OTPs are already compromised."
    ], "highlight")

    add_card(s2, Inches(6.86), Inches(4.5), col_w, row_h, "❓ Black-Box Confusion", [
        "Legacy security apps show cryptic risk percentages without explaining why a link is harmful.",
        "Non-technical users receive no clear, actionable instructions on what immediate steps to take."
    ], "highlight")

    add_speaker_notes(s2, "Phishing has moved to mobile messaging. Traditional security tools fail because they expect the user to be suspicious first and copy-paste links. But social engineering relies on panic. We need a system that detects threats autonomously the millisecond they arrive.")

    # =========================================================================
    # SLIDE 3: Proposed Solution
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "Core Solution", "💡 CyberShield: Autonomous Zero-Friction Defense", "Proactive background intelligence bridging threat arrival and user action", 3)

    t_w = Inches(3.64)
    t_h = Inches(4.75)
    add_card(s3, Inches(0.8), Inches(2.0), t_w, t_h, "⚡ Background Auto-Intercept", [
        "Native Android service silently monitors incoming WhatsApp, Telegram & SMS notifications.",
        "Zero user friction: scans run in the background without needing the app to be open.",
        "Instant high-priority lockscreen alerts with custom vibration patterns for threat vectors."
    ], "highlight")

    add_card(s3, Inches(4.84), Inches(2.0), t_w, t_h, "🧠 Dual-Engine Hybrid AI", [
        "70% Deterministic Rules: Catches typosquatting, DGA domains, blacklisted TLDs, and OTP traps.",
        "30% NLP Machine Learning: TF-IDF vectorizer identifying semantic coercion and phishing syntax.",
        "Balances accuracy with strict false-positive prevention."
    ], "success")

    add_card(s3, Inches(8.88), Inches(2.0), t_w, t_h, "🔍 Explainable AI (XAI)", [
        "Clear color-coded verdicts: SAFE (Green), SUSPICIOUS (Amber), DANGEROUS (Red).",
        "Itemizes exact Red Flag triggers (e.g. Typosquatting, Missing HTTPS).",
        "Provides plain-language explanations and step-by-step actionable safety tips."
    ], "highlight")

    add_speaker_notes(s3, "CyberShield transforms mobile security from reactive to autonomous. It runs quietly in the background, intercepts incoming messages, executes hybrid AI inference in under 10 milliseconds, and warns you immediately if a threat is detected.")

    # =========================================================================
    # SLIDE 4: Key Features & USP
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Product Features", "🌟 Key Features & Unique Selling Propositions", "Comprehensive multi-vector threat protection across digital ecosystems", 4)

    add_card(s4, Inches(0.8), Inches(2.0), col_w, row_h, "💬 WhatsApp & Messaging Auto-Shield", [
        "Monitors WhatsApp, WhatsApp Business, Telegram, SMS, Gmail, and Instagram.",
        "Configured with VISIBILITY_PUBLIC to ensure threat details are never masked by Android's 'Content Hidden' lockscreen policy."
    ])

    add_card(s4, Inches(6.86), Inches(2.0), col_w, row_h, "🌐 Deep URL Security Heuristics", [
        "Real-time live DNS resolution verifying active domains vs sinkholes.",
        "Catches deceptive brand typosquatting (e.g. paypal-verify.xyz, sbi-reward.online).",
        "Detects IP-based hosts, excessive subdomains, and URL shortener abuse."
    ])

    add_card(s4, Inches(0.8), Inches(4.5), col_w, row_h, "📷 Smart Hardware QR Scanner", [
        "Direct camera hardware stream with flash toggle, camera flip, and gallery image decoding.",
        "Intelligently auto-routes between website URLs and raw text payloads without throwing exceptions."
    ])

    add_card(s4, Inches(6.86), Inches(4.5), col_w, row_h, "📋 1-Tap Clipboard Quick Check & History", [
        "Scan copied links or text snippets straight from the system clipboard with a single tap.",
        "SQLite-backed persistent history log enabling incident review and audit trails."
    ])

    add_speaker_notes(s4, "CyberShield covers every threat vector — whether it's an incoming text, a banking link, a compromised QR code at a store, or a link copied from social media. Everything is analyzed with uniform intelligence.")

    # =========================================================================
    # SLIDE 5: Architecture
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Technical Architecture", "🏗️ System Architecture & Data Pipeline", "High-throughput, asynchronous client-server architecture", 5)

    arch_w = Inches(5.66)
    arch_h = Inches(4.75)
    add_card(s5, Inches(0.8), Inches(2.0), arch_w, arch_h, "📱 Frontend & Client Layer", [
        "Flutter Mobile Client: Built with Provider state management, Dio HTTP networking with dynamic IP interceptor, dark glassmorphic UI.",
        "Android Native Bridge: Kotlin NotificationListenerService and SmsReceiver coupled via MethodChannel and EventChannel.",
        "React 19 Web Dashboard: Vite 6, Tailwind CSS, drag-and-drop QR file upload dropzone.",
        "Responsive cross-platform parity sharing the same centralized AI backend."
    ], "highlight")

    add_card(s5, Inches(6.86), Inches(2.0), arch_w, arch_h, "⚙️ Backend AI Services (FastAPI)", [
        "REST Gateway (:8000): Asynchronous endpoints for /api/analyze/url, /api/analyze/text, /api/analyze/qr, /api/history.",
        "Analysis Pipeline: Deterministic Rule Engine + TF-IDF ML Model + Explainable AI module.",
        "Database Layer: SQLite with WAL mode for persistent scan telemetry and audit logging.",
        "Average inference turnaround in sub-10 milliseconds."
    ], "success")

    add_speaker_notes(s5, "Here is our architecture. Our mobile app leverages native Android background services to capture notification streams. These are dispatched to our high-throughput FastAPI backend over asynchronous REST APIs, evaluated against our hybrid detection core, and returned within milliseconds.")

    # =========================================================================
    # SLIDE 6: AI Engine
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "AI & Heuristics", "🧠 The Hybrid AI Detection Engine", "Combining deterministic cybersecurity heuristics with NLP statistical models", 6)

    add_card(s6, Inches(0.8), Inches(2.0), arch_w, arch_h, "📐 Rule Engine (70% Deterministic Weight)", [
        "Typosquatting: Detects deceptive brand imitations across 20+ top services (Google, PayPal, SBI, Amazon, Netflix).",
        "DGA Detection: Calculates vowel ratios and consonant clustering to flag randomly generated malware domains.",
        "Urgency Heuristics: Identifies coercive phrases ('Account suspended', 'Immediate action required').",
        "Safe Context Engine: Whitelists legitimate transactional OTP formats ('Your OTP is 482910') to avoid false alarms."
    ], "highlight")

    add_card(s6, Inches(6.86), Inches(2.0), arch_w, arch_h, "🤖 Machine Learning Classifier (30% Weight)", [
        "Trained on 800+ real-world phishing and legitimate communications.",
        "TF-IDF N-gram feature extraction capturing semantic phishing cues.",
        "Balanced Verdict Logic:",
        "  • Dangerous: combined >= 48 or rule_score >= 50",
        "  • Suspicious: rule_score >= 15 or combined >= 18",
        "  • Safe: Clean input (0 rules) -> 85-95% Safe confidence score."
    ], "success")

    add_speaker_notes(s6, "Pure machine learning models often hallucinate or produce false alarms, while pure rules miss novel phrasing. CyberShield combines the interpretability and precision of deterministic security rules with the flexibility of ML NLP, delivering dependable verdicts.")

    # =========================================================================
    # SLIDE 7: Android Integration
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "Mobile Engineering", "📱 Deep Android OS Integration & Resilience", "Native Kotlin services engineered for real-world background execution", 7)

    add_card(s7, Inches(0.8), Inches(2.0), col_w, row_h, "🔔 Dual-Channel Notification Architecture", [
        "Threat Channel (High Importance): Fires with custom vibration patterns (500ms bursts) and red alert banners for Dangerous/Suspicious threats.",
        "Safe Channel (Default Importance): Provides non-intrusive confirmation that incoming messages were scanned and verified safe."
    ])

    add_card(s7, Inches(6.86), Inches(2.0), col_w, row_h, "🔓 Zero-Masking Lockscreen Visibility", [
        "Explicitly configured with Notification.VISIBILITY_PUBLIC.",
        "Eliminates Android 12/13/14's default 'Content Hidden' lockscreen masking so threat details are immediately visible."
    ])

    add_card(s7, Inches(0.8), Inches(4.5), col_w, row_h, "⚙️ Dynamic Network IP Synchronization", [
        "Native services read backend server IP dynamically from SharedPreferences (flutter.server_ip).",
        "Allows the app to seamlessly adapt when testing across changing WiFi networks without re-compilation."
    ], "highlight")

    add_card(s7, Inches(6.86), Inches(4.5), col_w, row_h, "🔋 Battery & Memory Efficiency", [
        "Event-driven architecture: zero background CPU polling when idle.",
        "Executes analysis on dedicated background worker threads without blocking main UI rendering."
    ], "highlight")

    add_speaker_notes(s7, "Our Android implementation is built for real-world resilience. We solved platform-specific challenges like background process limits, notification permissions, and lockscreen privacy masking to ensure users get immediate, clear visibility of potential threats.")

    # =========================================================================
    # SLIDE 8: Web Dashboard
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Web Platform", "💻 React 19 Web Command Center", "Modern desktop security suite for power users and administrators", 8)

    add_card(s8, Inches(0.8), Inches(2.0), col_w, row_h, "🖥️ Unified Multi-Vector Scanner", [
        "Dedicated interactive tabs for URL Scans, Email / SMS Text Scans, and QR Code Files.",
        "Instant visual validation preventing invalid inputs or empty queries."
    ])

    add_card(s8, Inches(6.86), Inches(2.0), col_w, row_h, "📂 Drag-and-Drop QR Dropzone", [
        "HTML5 drag-and-drop file upload with client-side image preview.",
        "Decodes QR images instantly via OpenCV / Pyzbar backend processing."
    ])

    add_card(s8, Inches(0.8), Inches(4.5), col_w, row_h, "🌐 Live DNS Verification", [
        "Performs active DNS lookups to catch unregistered domains, parked ad-farms, and sinkholed malicious links."
    ], "highlight")

    add_card(s8, Inches(6.86), Inches(4.5), col_w, row_h, "📊 Real-Time SQLite Audit Feed", [
        "Displays a chronological table of all past scans, confidence scores, and triggered rules.",
        "Built-in Incident Reporting feature allowing users to flag false positives for model retraining."
    ], "highlight")

    add_speaker_notes(s8, "For desktop users and enterprise administrators, our Web Dashboard provides a fast, modern command center. It shares the exact same backend engine as our mobile app, providing consistent threat intelligence across devices.")

    # =========================================================================
    # SLIDE 9: Explainable AI
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "User Experience & XAI", "🔍 Explainable AI (XAI) & Education", "Empowering users with transparent reasoning rather than black-box scores", 9)

    add_card(s9, Inches(0.8), Inches(2.0), col_w, row_h, "🏷️ Itemized Red Flag Breakdown", [
        "Itemizes exact security violations (e.g. Typosquatting Detected, Missing HTTPS Protocol, Urgency Language).",
        "Assigns severity tags (High, Medium, Low) for transparent risk evaluation."
    ], "danger")

    add_card(s9, Inches(6.86), Inches(2.0), col_w, row_h, "📈 Statistical Confidence Gauge", [
        "Displays quantified percentage agreement between ML classifier and deterministic rule sets.",
        "Helps users gauge the certainty level of the analysis."
    ], "highlight")

    add_card(s9, Inches(0.8), Inches(4.5), col_w, row_h, "💡 Actionable Safety Guidelines", [
        "Provides practical, step-by-step instructions (e.g. 'Do not share OTP', 'Navigate to the official portal manually').",
        "Transforms security alerts into proactive user education."
    ], "success")

    add_card(s9, Inches(6.86), Inches(4.5), col_w, row_h, "🎨 Slide-Up Result Bottom Sheet", [
        "Animated modal with bold color badges (Safe, Suspicious, Dangerous).",
        "One-tap dismiss with automatic camera resumption for subsequent scans."
    ])

    add_speaker_notes(s9, "Security tools shouldn't just block; they should educate. When CyberShield detects a threat, it doesn't just say 'Blocked'. It shows the exact red flags found and tells the user exactly what to do next. This builds long-term user awareness.")

    # =========================================================================
    # SLIDE 10: Metrics
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_header(s10, "Validation & Benchmarks", "📊 Testing, Verification & Performance Metrics", "Rigorously tested across real-world phishing datasets and edge cases", 10)

    add_card(s10, Inches(0.8), Inches(2.0), col_w, row_h, "✅ 28 / 28 Automated Tests Passed", [
        "100% test pass rate across unit, integration, and API test suites.",
        "Full coverage for URL normalizer, rate limiter, QR decoder, and false-positive reduction."
    ], "success")

    add_card(s10, Inches(6.86), Inches(2.0), col_w, row_h, "⚡ Sub-10 Millisecond Latency", [
        "Average backend inference turnaround in <10ms.",
        "Ensures zero lag on incoming mobile message notifications."
    ], "highlight")

    add_card(s10, Inches(0.8), Inches(4.5), col_w, row_h, "🎯 Zero False Positives on Clean Text", [
        "Casual text ('Hi bro let us meet tomorrow') scores 88.5% Safe with 0 flags.",
        "Legitimate bank transactional OTP messages correctly recognized as safe."
    ])

    add_card(s10, Inches(6.86), Inches(4.5), col_w, row_h, "🚨 96%+ Precision on Malicious URLs", [
        "Impersonated domains (paypal-security-update.xyz) reliably flagged with 96.0% Dangerous confidence."
    ], "danger")

    add_speaker_notes(s10, "We benchmarked and verified our entire stack with 28 automated integration tests covering URL parsing, false-positive prevention, OTP safety, and QR decoders. Everything executes in sub-10 milliseconds with high accuracy.")

    # =========================================================================
    # SLIDE 11: Roadmap
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)
    add_header(s11, "Future Vision", "🚀 Scaling CyberShield: Future Roadmap", "Expanding the ecosystem across platforms and emerging threat vectors", 11)

    add_card(s11, Inches(0.8), Inches(2.0), col_w, row_h, "🧩 Browser Extension Ecosystem", [
        "Chromium & Firefox extension providing live DOM phishing link scanning and visual threat badges."
    ])

    add_card(s11, Inches(6.86), Inches(2.0), col_w, row_h, "🗣️ Vernacular & Multilingual AI", [
        "Expanding NLP models to detect smishing in Indian regional languages (Telugu, Hindi, Tamil, Kannada).",
        "Catching phonetically spelled financial scam lures (Hinglish / Telugish)."
    ], "highlight")

    add_card(s11, Inches(0.8), Inches(4.5), col_w, row_h, "🔓 Dark Web & Identity Shield", [
        "Integration with breach databases to alert users if their emails or passwords appear in active leaks."
    ])

    add_card(s11, Inches(6.86), Inches(4.5), col_w, row_h, "🏢 Enterprise SIEM & SOC Connectors", [
        "Exporting mobile threat telemetry to corporate Splunk and Elastic clusters for enterprise fleet protection."
    ])

    add_speaker_notes(s11, "Looking ahead, our vision is to expand CyberShield into browser extensions, integrate vernacular language processing for regional scams, and add identity breach monitoring to create an all-in-one digital defense platform.")

    # =========================================================================
    # SLIDE 12: Conclusion
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12)

    top_bar12 = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
    top_bar12.fill.solid()
    top_bar12.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar12.line.color.rgb = COLOR_PRIMARY

    # Hero Shield Icon
    icon12 = s12.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    p_i12 = icon12.text_frame.paragraphs[0]
    p_i12.text = "🛡️"
    p_i12.font.size = Pt(44)
    p_i12.alignment = PP_ALIGN.CENTER

    # Title
    t12 = s12.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(0.8))
    p_t12 = t12.text_frame.paragraphs[0]
    p_t12.text = "CYBERSHIELD"
    p_t12.font.size = Pt(36)
    p_t12.font.bold = True
    p_t12.font.color.rgb = COLOR_TEXT_MAIN
    p_t12.font.name = "Arial"
    p_t12.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub12 = s12.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.333), Inches(0.5))
    p_sub12 = sub12.text_frame.paragraphs[0]
    p_sub12.text = "Protecting Every Click. Securing Every Message."
    p_sub12.font.size = Pt(15)
    p_sub12.font.bold = True
    p_sub12.font.color.rgb = COLOR_SECONDARY
    p_sub12.font.name = "Arial"
    p_sub12.alignment = PP_ALIGN.CENTER

    # 3 Summary Pillar Cards
    c12_w = Inches(3.6)
    c12_h = Inches(1.5)
    c12_y = Inches(3.3)

    pc1 = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), c12_y, c12_w, c12_h)
    pc1.fill.solid()
    pc1.fill.fore_color.rgb = COLOR_CARD_BG
    pc1.line.color.rgb = COLOR_SUCCESS
    pc1.text_frame.margin_top = Inches(0.2)
    pp1 = pc1.text_frame.paragraphs[0]
    pp1.text = "⚡ AUTONOMOUS DEFENSE"
    pp1.font.size = Pt(11)
    pp1.font.bold = True
    pp1.font.color.rgb = COLOR_SUCCESS
    pp1.alignment = PP_ALIGN.CENTER
    pp1_s = pc1.text_frame.add_paragraph()
    pp1_s.text = "Zero-friction background scanning on WhatsApp & SMS without user intervention."
    pp1_s.font.size = Pt(10)
    pp1_s.font.color.rgb = COLOR_TEXT_MAIN
    pp1_s.alignment = PP_ALIGN.CENTER

    pc2 = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.86), c12_y, c12_w, c12_h)
    pc2.fill.solid()
    pc2.fill.fore_color.rgb = COLOR_CARD_BG
    pc2.line.color.rgb = COLOR_PRIMARY
    pc2.text_frame.margin_top = Inches(0.2)
    pp2 = pc2.text_frame.paragraphs[0]
    pp2.text = "🧠 EXPLAINABLE AI (XAI)"
    pp2.font.size = Pt(11)
    pp2.font.bold = True
    pp2.font.color.rgb = COLOR_PRIMARY
    pp2.alignment = PP_ALIGN.CENTER
    pp2_s = pc2.text_frame.add_paragraph()
    pp2_s.text = "Itemized Red Flag security violations with actionable safety guidance."
    pp2_s.font.size = Pt(10)
    pp2_s.font.color.rgb = COLOR_TEXT_MAIN
    pp2_s.alignment = PP_ALIGN.CENTER

    pc3 = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.73), c12_y, c12_w, c12_h)
    pc3.fill.solid()
    pc3.fill.fore_color.rgb = COLOR_CARD_BG
    pc3.line.color.rgb = COLOR_SECONDARY
    pc3.text_frame.margin_top = Inches(0.2)
    pp3 = pc3.text_frame.paragraphs[0]
    pp3.text = "📊 PRODUCTION-READY"
    pp3.font.size = Pt(11)
    pp3.font.bold = True
    pp3.font.color.rgb = COLOR_SECONDARY
    pp3.alignment = PP_ALIGN.CENTER
    pp3_s = pc3.text_frame.add_paragraph()
    pp3_s.text = "28/28 automated tests passed with sub-10ms inference turnaround."
    pp3_s.font.size = Pt(10)
    pp3_s.font.color.rgb = COLOR_TEXT_MAIN
    pp3_s.alignment = PP_ALIGN.CENTER

    # Thank You Banner
    thanks_box = s12.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.333), Inches(0.8))
    p_th = thanks_box.text_frame.paragraphs[0]
    p_th.text = "Thank You! We are open for Questions."
    p_th.font.size = Pt(20)
    p_th.font.bold = True
    p_th.font.color.rgb = COLOR_TEXT_MAIN
    p_th.font.name = "Arial"
    p_th.alignment = PP_ALIGN.CENTER

    add_speaker_notes(s12, "To conclude, CyberShield bridges the critical gap between threat arrival and user action through intelligent, real-time background protection. Thank you judges. We are now open for questions!")

    # Save presentation
    prs.save(PPTX_PATH)
    print(f"Successfully generated PowerPoint PPTX at: {PPTX_PATH}")

if __name__ == "__main__":
    build_presentation()
